"""PPT生成器

生成演示文稿：
- 课程PPT
- 学术汇报
- 培训材料
- 专题讲座
"""
import asyncio
import logging
from typing import Dict, Any, List, Optional
from datetime import datetime
import os
import json

from .base import BaseGenerator, GenerationRequest, GenerationResult, OutputFormat

logger = logging.getLogger(__name__)


class PPTGenerator(BaseGenerator):
    """PPT生成器"""

    def __init__(self, output_dir: str = "data/outputs/ppt"):
        super().__init__()
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def validate_request(self, request: GenerationRequest) -> bool:
        """验证请求参数"""
        if not request.topic:
            return False

        slide_count = request.parameters.get("slide_count", 10)
        if slide_count < 1 or slide_count > 100:
            logger.warning(f"幻灯片数量超出范围: {slide_count}")
            return False

        return True

    async def generate(self, request: GenerationRequest) -> GenerationResult:
        """生成PPT"""
        try:
            self.logger.info(f"开始生成PPT: {request.topic}")

            # 获取PPT参数
            slide_count = request.parameters.get("slide_count", 10)
            style = request.parameters.get("style", "academic")
            theme = request.parameters.get("theme", "default")
            language = request.parameters.get("language", "zh")

            # 构建PPT结构
            ppt_structure = await self._build_ppt_structure(
                topic=request.topic,
                slide_count=slide_count,
                style=style,
                language=language
            )

            # 生成PPT文件
            output_path = await self._create_ppt(
                task_id=request.task_id,
                structure=ppt_structure,
                theme=theme
            )

            return GenerationResult(
                task_id=request.task_id,
                status=GenerationStatus.COMPLETED,
                output_path=output_path,
                output_url=f"/outputs/ppt/{os.path.basename(output_path)}",
                metadata={
                    "slide_count": len(ppt_structure['slides']),
                    "style": style,
                    "theme": theme
                },
                completed_at=datetime.now()
            )

        except Exception as e:
            logger.error(f"PPT生成失败: {e}", exc_info=True)
            return GenerationResult(
                task_id=request.task_id,
                status=GenerationStatus.FAILED,
                error_message=str(e)
            )

    async def _build_ppt_structure(
        self,
        topic: str,
        slide_count: int,
        style: str,
        language: str
    ) -> Dict[str, Any]:
        """构建PPT结构"""

        # 从知识库检索内容
        from services.retrieval.vector import VectorRetrievalService
        retrieval_service = VectorRetrievalService()

        # 检索主题相关内容
        search_results = await retrieval_service.search(
            query=topic,
            limit=slide_count
        )

        # 构建幻灯片
        slides = []

        # 封面页
        slides.append({
            "type": "title",
            "title": topic,
            "subtitle": f"灵知知识系统 - {datetime.now().strftime('%Y年%m月')}",
            "layout": "title_slide"
        })

        # 目录页
        if slide_count > 5:
            slides.append({
                "type": "toc",
                "title": "目录",
                "content": self._generate_topic_outline(search_results, slide_count),
                "layout": "toc_slide"
            })

        # 内容页
        for i, result in enumerate(search_results[:slide_count-2], 1):
            content = result.get('content', '')
            title = result.get('metadata', {}).get('title', f'第{i}部分')

            # 提取要点
            bullet_points = self._extract_bullet_points(content)

            slides.append({
                "type": "content",
                "title": title,
                "content": bullet_points,
                "notes": content[:200],  # 演讲者备注
                "layout": "content_slide"
            })

        # 总结页
        slides.append({
            "type": "summary",
            "title": "总结",
            "content": [
                "回顾核心要点",
                "强调实践意义",
                "鼓励深入学习"
            ],
            "layout": "content_slide"
        })

        return {
            "title": topic,
            "author": "灵知知识系统",
            "created": datetime.now().isoformat(),
            "slides": slides
        }

    def _generate_topic_outline(self, search_results: List[Dict], max_items: int) -> List[str]:
        """生成主题大纲"""
        outline = []
        for i, result in enumerate(search_results[:max_items-2], 1):
            title = result.get('metadata', {}).get('title', f'第{i}部分')
            outline.append(title)
        return outline

    def _extract_bullet_points(self, content: str, max_points: int = 5) -> List[str]:
        """从内容中提取要点"""
        # 简单实现：按句号分割
        sentences = content.split('。')
        bullet_points = []

        for sentence in sentences[:max_points]:
            sentence = sentence.strip()
            if len(sentence) > 5:
                bullet_points.append(sentence)

        # 如果没有提取到足够的要点，使用默认
        if not bullet_points:
            bullet_points = [
                "核心概念解析",
                "理论与实践结合",
                "重点要点说明",
                "注意事项提醒"
            ]

        return bullet_points

    async def _create_ppt(
        self,
        task_id: str,
        structure: Dict[str, Any],
        theme: str
    ) -> str:
        """创建PPT文件"""

        # 保存为JSON格式（可以被前端转换为PPT）
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"ppt_{task_id}_{timestamp}.json"
        filepath = os.path.join(self.output_dir, filename)

        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(structure, f, ensure_ascii=False, indent=2)

        logger.info(f"PPT结构已保存: {filepath}")

        # 如果安装了python-pptx，可以生成实际的PPTX文件
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            pptx_path = filepath.replace('.json', '.pptx')
            prs = Presentation()

            # 设置幻灯片尺寸（16:9）
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            for slide_data in structure['slides']:
                # 根据类型选择布局
                if slide_data['type'] == 'title':
                    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
                else:
                    slide_layout = prs.slide_layouts[1]  # 标题和内容

                slide = prs.slides.add_slide(slide_layout)

                # 设置标题
                title = slide.shapes.title
                title.text = slide_data['title']

                # 设置内容
                if slide_data['type'] != 'title' and len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()

                    for point in slide_data.get('content', []):
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0

            prs.save(pptx_path)
            logger.info(f"PPT文件已生成: {pptx_path}")

            return pptx_path

        except ImportError:
            logger.warning("python-pptx未安装，仅返回JSON格式")
            return filepath

        except Exception as e:
            logger.error(f"生成PPTX文件失败: {e}")
            return filepath

    async def generate_course_ppt(
        self,
        course_title: str,
        chapters: List[str],
        style: str = "teaching"
    ) -> str:
        """生成课程PPT"""

        structure = {
            "title": course_title,
            "type": "course",
            "slides": []
        }

        # 封面
        structure['slides'].append({
            "type": "title",
            "title": course_title,
            "subtitle": f"共{len(chapters)}章",
            "layout": "title_slide"
        })

        # 课程大纲
        structure['slides'].append({
            "type": "toc",
            "title": "课程大纲",
            "content": [f"第{i}章: {chapter}" for i, chapter in enumerate(chapters, 1)],
            "layout": "toc_slide"
        })

        # 为每章生成幻灯片
        for chapter in chapters:
            # 检索章节内容
            from services.retrieval.vector import VectorRetrievalService
            retrieval_service = VectorRetrievalService()

            search_results = await retrieval_service.search(
                query=f"{course_title} {chapter}",
                limit=3
            )

            for result in search_results:
                structure['slides'].append({
                    "type": "content",
                    "title": chapter,
                    "content": self._extract_bullet_points(result.get('content', '')),
                    "layout": "content_slide"
                })

        # 保存
        task_id = self._generate_task_id()
        return await self._create_ppt(task_id, structure, theme="default")
